# main.py
import os
import io
import tempfile
import subprocess
import logging
from pathlib import Path
from typing import List

import streamlit as st

# ---------- Logging ----------
logger = logging.getLogger("multimodal_app")
if not logger.handlers:
    ch = logging.StreamHandler()
    ch.setFormatter(logging.Formatter("%(asctime)s - %(levelname)s - %(message)s"))
    logger.addHandler(ch)
logger.setLevel(logging.INFO)

# ---------- Streamlit page config ----------
st.set_page_config(page_title="Multimodal Data Processor", layout="centered")
st.title("🎧 Multimodal Data Processor")
st.write("Upload documents, images, audio/video or paste a YouTube link. Ask natural language questions about the content.")

# ---------- Helpers ----------
def clean_text(text: str) -> str:
    if not text:
        return ""
    txt = " ".join(text.split())
    # keep printable characters
    txt = "".join(c for c in txt if 32 <= ord(c) <= 126 or c in "\n\r\t")
    return txt.strip()

def format_paragraph_preview(text: str, max_chars: int = 800) -> str:
    txt = text.strip()
    if len(txt) <= max_chars:
        return txt
    preview = txt[:max_chars].rsplit(".", 1)[0] + "."
    paragraphs = [p.strip() for p in preview.split(". ") if p.strip()]
    return "\n\n".join(p + "." for p in paragraphs)

# ---------- Lazy-load resources (cached) ----------
@st.cache_resource(show_spinner=False)
def load_embedding_model():
    try:
        from sentence_transformers import SentenceTransformer
    except Exception as e:
        st.error("Missing sentence-transformers. Install via requirements.txt")
        raise e
    return SentenceTransformer("all-MiniLM-L6-v2")

@st.cache_resource(show_spinner=False)
def load_whisper_model(model_name="tiny"):
    try:
        import whisper
        return ("whisper", whisper.load_model(model_name))
    except Exception as e:
        st.error("Whisper not installed. Install openai-whisper.")
        raise e


@st.cache_resource(show_spinner=False)
def init_faiss_index(dimension=384):
    try:
        import faiss
        import numpy as np
    except Exception as e:
        st.error("faiss-cpu not installed. Install faiss-cpu.")
        raise e
    index = faiss.IndexFlatL2(dimension)
    return index

# ---------- Simple in-memory vector store ----------
VECTOR_INDEX = None
STORED_CHUNKS: List[str] = []

def ensure_index(dimension=384):
    global VECTOR_INDEX, STORED_CHUNKS
    if VECTOR_INDEX is None:
        VECTOR_INDEX = init_faiss_index(dimension)
        STORED_CHUNKS = []
    return VECTOR_INDEX

def add_embeddings_to_index(embeddings, chunks):
    global VECTOR_INDEX, STORED_CHUNKS
    VECTOR_INDEX = ensure_index()
    import numpy as np
    emb_arr = np.array(embeddings).astype("float32")
    if emb_arr.ndim == 1:
        emb_arr = emb_arr.reshape(1, -1)
    VECTOR_INDEX.add(emb_arr)
    STORED_CHUNKS.extend(chunks)

def search_index_by_embedding(query_embedding, top_k=3):
    global VECTOR_INDEX, STORED_CHUNKS
    VECTOR_INDEX = ensure_index()
    import numpy as np
    q_arr = np.array([query_embedding]).astype("float32")
    D, I = VECTOR_INDEX.search(q_arr, top_k)
    results = []
    for idx in I[0]:
        if idx == -1:
            continue
        if 0 <= idx < len(STORED_CHUNKS):
            results.append(STORED_CHUNKS[idx])
    return results

# ---------- Text extraction (documents) ----------
def extract_text_from_bytes(file_bytes: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    text = ""
    try:
        if ext == ".pdf":
            # try pdfminer first
            try:
                from pdfminer.high_level import extract_text
                stream = io.BytesIO(file_bytes)
                text = extract_text(stream)
            except Exception:
                # fallback to PyMuPDF
                import fitz
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                pages = []
                for p in doc:
                    pages.append(p.get_text("text"))
                text = "\n".join(pages)
        elif ext == ".docx":
            import docx
            stream = io.BytesIO(file_bytes)
            # python-docx expects a path-like, so write temporary
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
            tmp.write(file_bytes)
            tmp.close()
            doc = docx.Document(tmp.name)
            text = "\n".join([p.text for p in doc.paragraphs])
            os.unlink(tmp.name)
        elif ext == ".pptx":
            from pptx import Presentation
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
            tmp.write(file_bytes)
            tmp.close()
            prs = Presentation(tmp.name)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)
            text = "\n".join(slides)
            os.unlink(tmp.name)
        elif ext in [".txt", ".md"]:
            text = file_bytes.decode(errors="ignore")
        else:
            # Unknown document extension
            text = ""
    except Exception as e:
        logger.exception("Document extraction failed")
        text = ""
    return clean_text(text)

# ---------- Image OCR ----------
def extract_text_from_image_bytes(file_bytes: bytes) -> str:
    try:
        from PIL import Image
        import pytesseract
    except Exception as e:
        st.error("Pillow or pytesseract not installed.")
        raise e
    try:
        img = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img)
        return clean_text(text)
    except Exception as e:
        logger.exception("Image OCR failed")
        return ""

# ---------- Audio/Video/YouTube transcription ----------
def download_youtube_audio(url: str, out_dir: str) -> str:
    try:
        import yt_dlp
    except Exception:
        st.error("yt-dlp not installed.")
        raise
    outtmpl = os.path.join(out_dir, "yt_audio.%(ext)s")
    ydl_opts = {
        "format": "bestaudio/best",
        "outtmpl": outtmpl,
        "quiet": True,
        "no_warnings": True,
        "postprocessors": [{
            "key": "FFmpegExtractAudio",
            "preferredcodec": "mp3",
            "preferredquality": "192",
        }]
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    mp3 = os.path.join(out_dir, "yt_audio.mp3")
    return mp3

def convert_to_wav(input_path: str, output_path: str):
    # use ffmpeg to produce a WAV mono 16k
    cmd = [
        "ffmpeg",
        "-y",
        "-i", input_path,
        "-ar", "16000",
        "-ac", "1",
        output_path
    ]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return output_path

def transcribe_audio_file(path_or_bytes, whisper_size="tiny"):
    # path_or_bytes: either path string or bytes
    # returns text
    loader_type, model_obj = load_whisper_model(whisper_size)
    text = ""
    temp_file = None
    try:
        if isinstance(path_or_bytes, bytes):
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
            tmp.write(path_or_bytes)
            tmp.close()
            input_path = tmp.name
            temp_file = tmp.name
        else:
            input_path = str(path_or_bytes)
        # If file is not wav, convert to wav
        if not input_path.lower().endswith(".wav"):
            wav_tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".wav")
            wav_tmp.close()
            convert_to_wav(input_path, wav_tmp.name)
            input_path = wav_tmp.name
            if temp_file:
                # original temp (if mp3 bytes) kept for deletion later
                pass
        if loader_type == "faster":
            # faster_whisper usage
            model = model_obj
            segments, info = model.transcribe(input_path)
            parts = [seg.text for seg in segments]
            text = " ".join(parts)
        else:
            # whisper usage
            model = model_obj
            res = model.transcribe(input_path)
            text = res.get("text", "")
    except Exception as e:
        logger.exception("Transcription failed")
        text = ""
    finally:
        # cleanup temp files if any
        try:
            if temp_file and os.path.exists(temp_file):
                os.unlink(temp_file)
        except Exception:
            pass
    return clean_text(text)

# ---------- Chunking and embeddings ----------
def get_text_chunks(text: str, chunk_size: int = 1000, overlap: int = 100):
    text = text.replace("\n", " ").strip()
    chunks = []
    start = 0
    n = len(text)
    if n == 0:
        return []
    while start < n:
        end = min(start + chunk_size, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

def embed_text_chunks(chunks: List[str]):
    model = load_embedding_model()
    embeddings = model.encode(chunks, show_progress_bar=False)
    return embeddings

# ---------- Gemini query handling ----------
def configure_gemini():
    try:
        import google.generativeai as genai
    except Exception:
        st.error("google-generativeai not installed. Install via requirements.")
        raise
    api_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GENAI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY environment variable not set. Set it before running.")
    genai.configure(api_key=api_key)
    return genai

def answer_with_gemini(user_query: str, context_chunks: List[str]):
    genai = configure_gemini()
    
    prompt = (
        "You are a helpful assistant. Use the following context to answer the question.\n"
        "Only answer if the context clearly supports it. If the answer isn't present, say 'I don't know'.\n\n"
        "=== CONTEXT START ===\n"
    )

    for i, chunk in enumerate(context_chunks):
        snippet = chunk if len(chunk) < 1000 else chunk[:1000] + "..."
        prompt += f"[Snippet {i+1}]: {snippet}\n\n"

    prompt += "=== CONTEXT END ===\n\n"
    prompt += f"Question: {user_query}\nAnswer:"

    try:
        m = genai.GenerativeModel("gemini-2.5-flash-lite")
        response = m.generate_content(prompt)

        # Try multiple extraction options
        if hasattr(response, "text"):
            return response.text.strip()
        elif hasattr(response, "candidates"):
            cand = response.candidates[0]
            try:
                return cand.content[0].text.strip()
            except Exception:
                try:
                    return cand.content.parts[0].text.strip()
                except Exception:
                    return str(cand)
        else:
            return "No answer generated by Gemini."
    except Exception as e:
        logger.exception("Gemini API call failed")
        return f"Gemini error: {e}"


# ---------- Streamlit UI layout ----------
with st.sidebar:
    st.header("Settings")
    chunk_size = st.number_input("Chunk size (chars)", min_value=200, max_value=20000, value=1000, step=100)
    chunk_overlap = st.number_input("Chunk overlap (chars)", min_value=0, max_value=chunk_size//2, value=100, step=50)
    whisper_model_choice = st.selectbox("Whisper model (transcription)", options=["tiny", "base", "small"], index=0)
    st.markdown("---")
    st.markdown("**Env variables**:")
    st.code("GEMINI_API_KEY=your_key_here")
    st.markdown("Tesseract & FFmpeg must be installed on the system for OCR/audio features.")

uploaded = st.file_uploader("Upload file (pdf, docx, pptx, txt, md, png, jpg, mp3, mp4) OR paste YouTube link below", type=None)
yt_link = st.text_input("Or paste a YouTube URL (optional)")

# area for extracted text preview
extracted_text = ""
indexed = False

if uploaded is not None or yt_link:
    # handle uploaded file
    if uploaded is not None:
        file_bytes = uploaded.getvalue()
        filename = uploaded.name
        st.success(f"Uploaded: {filename}")
        ext = Path(filename).suffix.lower()
        if ext in [".pdf", ".docx", ".pptx", ".txt", ".md"]:
            extracted_text = extract_text_from_bytes(file_bytes, filename)
        elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
            extracted_text = extract_text_from_image_bytes(file_bytes)
        elif ext in [".mp3", ".wav", ".m4a", ".mp4", ".mov", ".aac"]:
            st.info("Transcribing audio/video (this may take time)...")
            # write bytes to temp file to pass to ffmpeg/whisper
            tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
            tmp.write(file_bytes)
            tmp.close()
            extracted_text = transcribe_audio_file(tmp.name, whisper_size=whisper_model_choice)
            os.unlink(tmp.name)
        else:
            st.warning("Unsupported uploaded file type.")
    else:
        # youtube link provided
        tmpd = tempfile.mkdtemp()
        try:
            st.info("Downloading audio from YouTube (yt-dlp)...")
            mp3_path = download_youtube_audio(yt_link, tmpd)
            st.info("Transcribing YouTube audio (this may take time)...")
            extracted_text = transcribe_audio_file(mp3_path, whisper_size=whisper_model_choice)
        except Exception as e:
            logger.exception("YouTube download/transcription failed")
            st.error("YouTube processing failed: " + str(e))
        finally:
            # cleanup
            try:
                for f in Path(tmpd).glob("*"):
                    f.unlink()
                Path(tmpd).rmdir()
            except Exception:
                pass

    if not extracted_text:
        st.warning("No text was extracted from the provided input.")
    else:
        st.markdown("### 📄 Extracted Text (preview)")
        st.text_area("Preview", format_paragraph_preview(extracted_text, max_chars=1200), height=300)

        with st.spinner("🔎 Chunking, embedding, and indexing extracted text..."):
            chunks = get_text_chunks(extracted_text, chunk_size=chunk_size, overlap=chunk_overlap)
            if not chunks:
                st.error("No chunks produced (empty text).")
            else:
                try:
                    emb_model = load_embedding_model()
                    embeddings = emb_model.encode(chunks, show_progress_bar=False)
                    add_embeddings_to_index(embeddings, chunks)
                    st.success(f"✅ Automatically indexed {len(chunks)} content chunks.")
                    indexed = True
                except Exception as e:
                    logger.exception("Embedding/indexing failed")
                    st.error(f"Embedding/indexing failed: {e}")

# Query and answer
st.markdown("---")
st.header("💬 Ask a question")
query = st.text_input("Enter a natural language question about the indexed content")

if st.button("Search & Answer"):
    if not query.strip():
        st.warning("Enter a question first.")
    else:
        try:
            emb_model = load_embedding_model()
            q_emb = emb_model.encode([query])[0]
            top_chunks = search_index_by_embedding(q_emb, top_k=5)

            if not top_chunks:
                st.info("No relevant content indexed.")
            else:
                st.info("Generating answer from Gemini (may take a few seconds)...")
                try:
                    answer = answer_with_gemini(query, top_chunks)
                    st.markdown("### 🤖 Answer")
                    st.write(answer)
                except RuntimeError as re:
                    st.error(str(re))
        except Exception as e:
            logger.exception("Search/answer failed")
            st.error(f"Search/Answer failed: {e}")


# Footer
st.markdown("---")
st.caption("Developer: Shaik Rahamtulla • Email: rahamtullawork@gmail.com")
